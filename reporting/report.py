from pptx import Presentation #python-pptx module
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import pandas as pd
import logging
import os

logging.basicConfig()
logger = logging.getLogger()
logger.setLevel(logging.INFO)

path = os.environ['ATHENA_SUMMARY_CSV']
head, tail = os.path.split(path)

df = pd.read_csv(tail)
prs = Presentation('SHIP_HealthCheck_v3.0_Template.pptx')

slides = [slide for slide in prs.slides]
shapes = []
for slide in slides:
    for shape in slide.shapes:
        shapes.append(shape)

def replace_text(replacements: dict, shapes: list):
    """Takes dict of {match: replacement, ... } and replaces all matches.
    Currently not implemented for charts or graphics.
    """
    for shape in shapes:
        for match, replacement in replacements.items():
            if shape.has_text_frame:
                if (shape.text.find(match)) != -1:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        paragraph.alignment = PP_ALIGN.CENTER
                        for run in paragraph.runs:
                            cur_text = run.text
                            new_text = cur_text.replace(str(match), str(replacement))
                            run.text = new_text
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if match in cell.text:
                            new_text = cell.text.replace(match, replacement)
                            cell.text_frame.paragraphs[0].text = new_text
                            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    return shapes

unique_accounts = df['account_uid'].unique()
unique_checks = df['check_id'].unique()

critical_severity_count = ((df['severity'] == 'critical') & (df['status'] == 'FAIL')).value_counts().get(True, 0)
high_severity_count = ((df['severity'] == 'high') & (df['status'] == 'FAIL')).value_counts().get(True, 0)
medium_severity_count = ((df['severity'] == 'medium') & (df['status'] == 'FAIL')).value_counts().get(True, 0)
low_severity_count = ((df['severity'] == 'low') & (df['status'] == 'FAIL')).value_counts().get(True, 0)
total_findings = critical_severity_count + high_severity_count + medium_severity_count + low_severity_count

iam_root_mfa_enabled = ((df['check_id'] == 'iam_root_mfa_enabled') & (df['status'] == 'FAIL')).value_counts().get(True, 0)
ec2_securitygroup_allow_ingress_from_internet_to_any_port = ((df['check_id'] == 'ec2_securitygroup_allow_ingress_from_internet_to_any_port') & (df['status'] == 'FAIL')).value_counts().get(True, 0)
s3_bucket_public_access = ((df['check_id'] == 's3_bucket_public_access') & (df['status'] == 'FAIL')).value_counts().get(True, 0)
cloudtrail_multi_region_enabled = ((df['check_id'] == 'cloudtrail_multi_region_enabled') & (df['status'] == 'FAIL')).value_counts().get(True, 0)
guardduty_is_enabled = ((df['check_id'] == 'guardduty_is_enabled') & (df['status'] == 'FAIL')).value_counts().get(True, 0)
iam_rotate_access_key_90_days = ((df['check_id'] == 'iam_rotate_access_key_90_days') & (df['status'] == 'FAIL')).value_counts().get(True, 0)
iam_password_policy_lowercase = ((df['check_id'] == 'iam_password_policy_lowercase') & (df['status'] == 'FAIL')).value_counts().get(True, 0)
iam_password_policy_number = ((df['check_id'] == 'iam_password_policy_number') & (df['status'] == 'FAIL')).value_counts().get(True, 0)
iam_password_policy_symbol = ((df['check_id'] == 'iam_password_policy_symbol') & (df['status'] == 'FAIL')).value_counts().get(True, 0)
iam_password_policy_uppercase = ((df['check_id'] == 'iam_password_policy_uppercase') & (df['status'] == 'FAIL')).value_counts().get(True, 0)
config_recorder_all_regions_enabled = ((df['check_id'] == 'config_recorder_all_regions_enabled') & (df['status'] == 'FAIL')).value_counts().get(True, 0)
awslambda_function_using_supported_runtimes = ((df['check_id'] == 'awslambda_function_using_supported_runtimes') & (df['status'] == 'FAIL')).value_counts().get(True, 0)

iam_password_policy = iam_password_policy_lowercase + iam_password_policy_number + iam_password_policy_symbol + iam_password_policy_uppercase

logger.info("Data summary")
logger.info(f"Critical: {critical_severity_count}")
logger.info(f"High: {high_severity_count}")
logger.info(f"Medium: {medium_severity_count}")
logger.info(f"Low: {low_severity_count}")
logger.info(f"Total: {total_findings}")

logger.info(f"Accounts analyzed: {len(unique_accounts)}")
logger.info(f"Checks performed: {len(unique_checks)}")
logger.info(f"MFA on Root User: {iam_root_mfa_enabled}")
logger.info(f"Unrestricted Security Groups: {ec2_securitygroup_allow_ingress_from_internet_to_any_port}")
logger.info(f"Public S3 Buckets: {s3_bucket_public_access}")
logger.info(f"Multi-Region CloudTrail: {cloudtrail_multi_region_enabled}")
logger.info(f"GuardDuty Enabled: {guardduty_is_enabled}")
logger.info(f"Access Key Rotation: {iam_rotate_access_key_90_days}")
logger.info(f"Password Policy (All findings): {iam_password_policy}")
logger.info(f"Config Recorder (All Regions): {config_recorder_all_regions_enabled}")
logger.info(f"Lambda Function (Supported Runtimes): {awslambda_function_using_supported_runtimes}")

replace_text({'[+unique_accounts+]': str(len(unique_accounts))}, shapes)
replace_text({'[+unique_checks+]': str(len(unique_checks))}, shapes)
replace_text({'[+critical_severity_count+]': str(critical_severity_count)}, shapes)
replace_text({'[+high_severity_count+]': str(high_severity_count)}, shapes)
replace_text({'[+medium_severity_count+]': str(medium_severity_count)}, shapes)
replace_text({'[+low_severity_count+]': str(low_severity_count)}, shapes)
replace_text({'[+total_findings+]': str(total_findings)}, shapes)
replace_text({'[+iam_root_mfa_enabled+]': str(iam_root_mfa_enabled)}, shapes)
replace_text({'[+ec2_securitygroup_allow_ingress_from_internet_to_any_port+]': str(ec2_securitygroup_allow_ingress_from_internet_to_any_port)}, shapes)
replace_text({'[+s3_bucket_public_access+]': str(s3_bucket_public_access)}, shapes)
replace_text({'[+cloudtrail_multi_region_enabled+]': str(cloudtrail_multi_region_enabled)}, shapes)
replace_text({'[+guardduty_is_enabled+]': str(guardduty_is_enabled)}, shapes)
replace_text({'[+iam_rotate_access_key_90_days+]': str(iam_rotate_access_key_90_days)}, shapes)
replace_text({'[+iam_password_policy+]': str(iam_password_policy)}, shapes)
replace_text({'[+config_recorder_all_regions_enabled+]': str(config_recorder_all_regions_enabled)}, shapes)
replace_text({'[+awslambda_function_using_supported_runtimes+]': str(awslambda_function_using_supported_runtimes)}, shapes)

prs.save('SHIP_v3_Updated.pptx')